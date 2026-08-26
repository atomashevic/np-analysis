"""Convert Obsidian slide markdown to PPTX for Google Slides upload."""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Config ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors — minimalist white
BG_DARK = RGBColor(0xFF, 0xFF, 0xFF)
BG_SECTION = RGBColor(0xF8, 0xF8, 0xF8)
ACCENT = RGBColor(0x1A, 0x73, 0xE8)  # Google blue
TEXT_WHITE = RGBColor(0x20, 0x20, 0x20)  # near-black
TEXT_LIGHT = RGBColor(0x3C, 0x40, 0x43)  # dark gray
TEXT_MUTED = RGBColor(0x80, 0x86, 0x8B)  # medium gray
TABLE_BG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEADER_BG = RGBColor(0xF1, 0xF3, 0xF4)
TABLE_ALT_BG = RGBColor(0xF8, 0xF9, 0xFA)

ASSETS_DIR = Path(__file__).resolve().parent.parent / "assets"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=TEXT_WHITE,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Google Sans",
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(
    slide,
    left,
    top,
    width,
    height,
    items,
    font_size=16,
    color=TEXT_LIGHT,
    bold_prefix=True,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = Pt(8)
        p.space_before = Pt(4)

        # Handle **bold** sections
        parts = re.split(r"(\*\*.*?\*\*)", item)
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                run = p.add_run()
                run.text = part[2:-2]
                run.font.size = Pt(font_size)
                run.font.color.rgb = ACCENT
                run.font.bold = True
                run.font.name = "Google Sans"
            elif part.startswith("*") and part.endswith("*"):
                run = p.add_run()
                run.text = part[1:-1]
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.italic = True
                run.font.name = "Google Sans"
            else:
                run = p.add_run()
                run.text = part
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.name = "Google Sans"
    return tf


def add_table(slide, left, top, width, rows_data, col_widths=None):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 0
    if n_rows == 0 or n_cols == 0:
        return

    row_height = Inches(0.4)
    table_height = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = cell_text.strip()

            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Google Sans"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = ACCENT
                else:
                    p.font.color.rgb = TEXT_LIGHT
                # Right-align numeric columns
                if c > 0 and r > 0:
                    p.alignment = PP_ALIGN.RIGHT
                elif r == 0:
                    p.alignment = PP_ALIGN.CENTER

            # Cell fill
            fill = cell.fill
            fill.solid()
            if r == 0:
                fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                fill.fore_color.rgb = TABLE_ALT_BG
            else:
                fill.fore_color.rgb = TABLE_BG


def add_accent_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_image_if_exists(slide, img_name, left, top, width, height=None):
    img_path = ASSETS_DIR / img_name
    if img_path.exists():
        if height:
            slide.shapes.add_picture(str(img_path), left, top, width, height)
        else:
            slide.shapes.add_picture(str(img_path), left, top, width=width)
        return True
    return False


def parse_table(lines):
    """Parse markdown table lines into list of rows."""
    rows = []
    for line in lines:
        line = line.strip()
        if line.startswith("|") and not re.match(r"^\|[\s\-:]+\|", line):
            cells = [c.strip() for c in line.split("|")[1:-1]]
            # Clean markdown bold
            cells = [re.sub(r"\*\*(.*?)\*\*", r"\1", c) for c in cells]
            rows.append(cells)
    return rows


def parse_slides(md_text):
    """Split markdown into slides by --- separator."""
    slides = re.split(r"\n---\n", md_text)
    return [s.strip() for s in slides if s.strip()]


def build_presentation(md_path, output_path):
    md_text = Path(md_path).read_text(encoding="utf-8")
    slide_texts = parse_slides(md_text)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    for idx, slide_text in enumerate(slide_texts):
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, BG_DARK)

        lines = slide_text.split("\n")

        # Extract title
        title = ""
        content_lines = []
        for line in lines:
            if line.startswith("# ") and not title:
                title = line[2:].strip()
            elif line.startswith("## ") and not title:
                title = line[3:].strip()
            else:
                content_lines.append(line)

        # --- Title slide (first) ---
        if idx == 0:
            set_slide_bg(slide, BG_SECTION)
            add_accent_bar(slide, Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(0.03))
            add_text_box(
                slide,
                Inches(1),
                Inches(1.5),
                Inches(11),
                Inches(1.5),
                title,
                font_size=36,
                color=TEXT_WHITE,
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            # Subtitle
            subtitle = ""
            for line in content_lines:
                if line.startswith("## "):
                    subtitle = line[3:].strip()
                elif (
                    line.startswith("*")
                    and line.endswith("*")
                    and not line.startswith("**")
                ):
                    subtitle += "\n" + line.strip("* ")
            if subtitle:
                add_text_box(
                    slide,
                    Inches(1),
                    Inches(3.2),
                    Inches(11),
                    Inches(1.5),
                    subtitle,
                    font_size=22,
                    color=TEXT_MUTED,
                    alignment=PP_ALIGN.CENTER,
                )
            # Slide number
            add_text_box(
                slide,
                Inches(12),
                Inches(7),
                Inches(1),
                Inches(0.4),
                str(idx + 1),
                font_size=10,
                color=TEXT_MUTED,
                alignment=PP_ALIGN.RIGHT,
            )
            continue

        # --- Regular slides ---
        # Title bar
        add_accent_bar(slide, Inches(0.5), Inches(0.55), Inches(0.04), Inches(0.45))
        add_text_box(
            slide,
            Inches(0.8),
            Inches(0.4),
            Inches(11),
            Inches(0.7),
            title,
            font_size=28,
            color=TEXT_WHITE,
            bold=True,
        )

        # Parse content
        bullets = []
        table_lines = []
        in_table = False
        in_code = False
        code_lines = []
        images = []
        extra_text = []

        for line in content_lines:
            stripped = line.strip()

            if stripped.startswith("```"):
                if in_code:
                    in_code = False
                    extra_text.append("\n".join(code_lines))
                    code_lines = []
                else:
                    in_code = True
                continue

            if in_code:
                code_lines.append(line)
                continue

            if stripped.startswith("|"):
                in_table = True
                table_lines.append(stripped)
            elif in_table and not stripped.startswith("|"):
                in_table = False
                if stripped.startswith("- ") or stripped.startswith("1. "):
                    bullets.append(re.sub(r"^[\-\d]+[\.\)]\s*", "", stripped))
                elif stripped.startswith("![["):
                    m = re.search(r"!\[\[(.*?)(?:\|.*)?\]\]", stripped)
                    if m:
                        images.append(m.group(1))
                elif stripped and not stripped.startswith("#"):
                    extra_text.append(stripped)
            elif stripped.startswith("- ") or stripped.startswith("  - "):
                bullets.append(re.sub(r"^\s*[\-]\s*", "", stripped))
            elif re.match(r"^\d+\.\s", stripped):
                bullets.append(re.sub(r"^\d+\.\s*", "", stripped))
            elif stripped.startswith("- [ ]"):
                bullets.append(re.sub(r"^- \[ \]\s*", "☐ ", stripped))
            elif stripped.startswith("![["):
                m = re.search(r"!\[\[(.*?)(?:\|.*)?\]\]", stripped)
                if m:
                    images.append(m.group(1))
            elif stripped and not stripped.startswith("#") and stripped != "":
                extra_text.append(stripped)

        # Determine layout
        has_table = len(table_lines) > 0
        has_image = len(images) > 0
        has_bullets = len(bullets) > 0

        content_top = Inches(1.3)
        content_left = Inches(0.8)
        content_width = Inches(11.5)

        if has_table and has_image:
            # Table + image side by side or stacked
            table_data = parse_table(table_lines)
            if table_data:
                add_table(slide, content_left, content_top, Inches(6.5), table_data)

            img_name = images[0].split("/")[-1] if "/" in images[0] else images[0]
            add_image_if_exists(slide, img_name, Inches(7.5), content_top, Inches(5.2))

            if bullets:
                bullet_top = content_top + Inches(len(table_data) * 0.4 + 0.3)
                add_bullet_list(
                    slide,
                    content_left,
                    bullet_top,
                    Inches(11.5),
                    Inches(2.5),
                    bullets,
                    font_size=15,
                )

        elif has_table and has_bullets:
            # Table on top, bullets below
            table_data = parse_table(table_lines)
            if table_data:
                add_table(slide, content_left, content_top, content_width, table_data)
                bullet_top = content_top + Inches(len(table_data) * 0.4 + 0.3)
            else:
                bullet_top = content_top

            add_bullet_list(
                slide,
                content_left,
                bullet_top,
                content_width,
                Inches(3),
                bullets,
                font_size=15,
            )

        elif has_table:
            table_data = parse_table(table_lines)
            if table_data:
                add_table(slide, content_left, content_top, content_width, table_data)
            if extra_text:
                et_top = content_top + Inches(len(table_data) * 0.4 + 0.3)
                for et in extra_text:
                    add_bullet_list(
                        slide,
                        content_left,
                        et_top,
                        content_width,
                        Inches(1.5),
                        [et],
                        font_size=15,
                    )
                    et_top += Inches(0.5)

        elif has_image and has_bullets:
            # Bullets left, image right
            add_bullet_list(
                slide,
                content_left,
                content_top,
                Inches(5.5),
                Inches(5),
                bullets,
                font_size=16,
            )

            img_name = images[0].split("/")[-1] if "/" in images[0] else images[0]
            add_image_if_exists(slide, img_name, Inches(6.8), content_top, Inches(6))

        elif has_image:
            # Full-width image
            img_name = images[0].split("/")[-1] if "/" in images[0] else images[0]
            if extra_text:
                add_bullet_list(
                    slide,
                    content_left,
                    content_top,
                    content_width,
                    Inches(1.5),
                    extra_text[:3],
                    font_size=15,
                )
                add_image_if_exists(
                    slide, img_name, Inches(1.5), Inches(2.5), Inches(10)
                )
            else:
                add_image_if_exists(slide, img_name, Inches(1), content_top, Inches(11))

        elif has_bullets:
            # Bullets only
            if extra_text:
                # Extra text as intro paragraph
                para_items = [et for et in extra_text if not et.startswith("```")]
                if para_items:
                    add_bullet_list(
                        slide,
                        content_left,
                        content_top,
                        content_width,
                        Inches(1.5),
                        para_items,
                        font_size=15,
                        color=TEXT_MUTED,
                    )
                    content_top += Inches(len(para_items) * 0.5 + 0.2)

            add_bullet_list(
                slide,
                content_left,
                content_top,
                content_width,
                Inches(5.5),
                bullets,
                font_size=17,
            )

        elif extra_text:
            # Code block or text only
            combined = "\n".join(extra_text)
            if len(combined) > 200:
                # Likely code block - use monospace
                txBox = slide.shapes.add_textbox(
                    content_left, content_top, content_width, Inches(5.5)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = combined
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_LIGHT
                p.font.name = "Roboto Mono"
            else:
                add_bullet_list(
                    slide,
                    content_left,
                    content_top,
                    content_width,
                    Inches(5.5),
                    extra_text,
                    font_size=17,
                )

        # Slide number
        add_text_box(
            slide,
            Inches(12),
            Inches(7),
            Inches(1),
            Inches(0.4),
            str(idx + 1),
            font_size=10,
            color=TEXT_MUTED,
            alignment=PP_ALIGN.RIGHT,
        )

    prs.save(output_path)
    print(f"Saved: {output_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    root = Path(__file__).resolve().parent.parent
    build_presentation(
        str(root / "nisamprijavila-slides.md"),
        str(root / "nisamprijavila_wip.pptx"),
    )
