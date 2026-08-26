#!/usr/bin/env python3
"""
Repair slide order in nisamprijavila_wip.pptx.

The add_institutional_slides.py used reverse move order which scrambled slides.
This script identifies slides by title fragment and reorders to the correct sequence.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_PATH = Path("/home/socio/nisamprijavila/nisamprijavila_wip.pptx")
TEXT_MUTED = RGBColor(0x80, 0x86, 0x8B)

# Desired title fragments in order (must be unique enough to match)
DESIRED_ORDER = [
    "From Solidarity to Broadcast",
    "Motivation",
    "The Case:",
    "Discursive Counterpublic",
    "Amplification Paradox",
    "Visibility, Erasure",
    "Comparative Cases",
    "Hypotheses",
    "Data",
    "Phase Detection",
    "Phase-Aggregated Network",
    "Early-Author Cohort Analysis",
    "Eight Detected Phases",
    "Reciprocity Collapse",
    "Attention Concentration",
    "Network Fragmentation",
    "Early-Author Advantage",
    "Phase 1 Cohort Activity",
    "Institutional vs. Regular",           # NEW
    "Phase 1 Authors vs. Institutions",    # NEW
    "Key Takeaway:",                        # NEW
    "Structural Arc",
    "Discussion:",
    "Contribution",
    "Planned Paper",
    "Next Steps",
    "References",
]


def get_slide_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text and not text.isdigit():
                return text[:80]
    return ""


def update_slide_numbers(prs):
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text.strip()
            if txt.isdigit() and shape.left >= int(11 * 914400) and shape.top >= int(6.5 * 914400):
                p = shape.text_frame.paragraphs[0]
                if p.runs:
                    p.runs[0].text = str(i + 1)
                else:
                    p.text = str(i + 1)
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT_MUTED
                p.font.name = "Google Sans"
                p.alignment = PP_ALIGN.RIGHT


def main():
    prs = Presentation(str(PPTX_PATH))

    print("Current order:")
    slides = list(prs.slides)
    for i, slide in enumerate(slides):
        print(f"  [{i:2d}] {get_slide_title(slide)}")

    # Match each desired fragment to current slide indices
    remaining = list(range(len(slides)))
    ordered_indices = []

    for frag in DESIRED_ORDER:
        found = None
        for pos, idx in enumerate(remaining):
            title = get_slide_title(slides[idx])
            if frag.lower() in title.lower():
                found = pos
                break
        if found is not None:
            ordered_indices.append(remaining.pop(found))
        else:
            print(f"  WARNING: no slide matching '{frag}'")

    if remaining:
        print(f"  WARNING: unmatched slides at indices {remaining}")
        ordered_indices.extend(remaining)

    print("\nDesired order (by matched index):")
    for new_pos, old_idx in enumerate(ordered_indices):
        print(f"  [{new_pos:2d}] <- [{old_idx:2d}] {get_slide_title(slides[old_idx])}")

    # Reorder sldIdLst in one shot
    sldIdLst = prs.slides._sldIdLst
    all_sldIds = list(sldIdLst)
    new_order = [all_sldIds[i] for i in ordered_indices]

    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for sldId in new_order:
        sldIdLst.append(sldId)

    update_slide_numbers(prs)
    prs.save(str(PPTX_PATH))
    print(f"\nSaved: {PPTX_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
